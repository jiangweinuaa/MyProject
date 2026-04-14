#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
智问核心流程 PPT 生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # ========== 封面页 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "智问系统核心流程"
    subtitle.text = "自然语言查询系统技术架构\n2026-04-14"
    
    # ========== 第 1 页：系统概述 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "系统概述"
    tf = content.text_frame
    tf.text = "智问是什么？"
    
    p = tf.add_paragraph()
    p.text = "用自然语言对话查询业务数据的智能系统"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "核心能力"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 理解自然语言问题（如今日销售额）"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• AI 自动生成 SQL 查询"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 安全验证后执行查询"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 返回结构化数据结果"
    p.level = 1
    
    # ========== 第 2 页：整体架构 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "整体架构"
    tf = content.text_frame
    tf.text = "四层架构"
    
    p = tf.add_paragraph()
    p.text = "1. 用户层：SmartQuery.vue 聊天界面"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "2. 前端层：REPORTUI（Vue 3 + Element Plus）"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. 后端层：REPORT Service（Spring Boot）"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "4. 数据层：Oracle 数据库 + 外部 AI 服务"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "API 端点：POST /api/nl-query/query"
    p.level = 1
    
    # ========== 第 3 页：核心流程 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "核心流程"
    tf = content.text_frame
    tf.text = "7 步完成一次查询"
    
    p = tf.add_paragraph()
    p.text = "1. 用户输入问题 → 前端发送 POST 请求"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "2. NLQueryController 接收请求"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. AISQLService 生成 SQL（调用通义千问 API）"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "4. SchemaService 提供表结构元数据"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "5. validateSQL 安全验证（白名单检查）"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "6. JdbcTemplate 执行 SQL 查询"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "7. 返回 JSON 结果 → 前端渲染表格"
    p.level = 1
    
    # ========== 第 4 页：AI 集成 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "AI 集成"
    tf = content.text_frame
    tf.text = "通义千问 API 调用"
    
    p = tf.add_paragraph()
    p.text = "模型：qwen-plus"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "API 端点：dashscope.aliyuncs.com"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "API Key：从 PRODUCT_APPKEY 表读取"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Prompt 构建"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 表结构（从 AI_TABLE_FILTER 获取）"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 用户问题"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 9 条 SQL 生成规范"
    p.level = 1
    
    # ========== 第 5 页：安全机制 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "安全机制"
    tf = content.text_frame
    tf.text = "三层安全验证"
    
    p = tf.add_paragraph()
    p.text = "1. SQL 类型检查"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   ✓ 只允许 SELECT 语句"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "   ✓ 禁止 DROP/DELETE/UPDATE/INSERT 等"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "2. 表白名单检查"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   ✓ 从 AI_TABLE_FILTER 表读取允许的表"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "   ✓ SQL 中必须包含允许的表名"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "3. 只读账号执行"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   ✓ JdbcTemplate 使用只读数据库账号"
    p.level = 2
    
    # ========== 第 6 页：核心数据表 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "核心数据表"
    tf = content.text_frame
    tf.text = "数据库表结构"
    
    p = tf.add_paragraph()
    p.text = "AI_TABLE_FILTER"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 控制 AI 可访问的表白名单"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 字段：TABLE_NAME, TABLE_COMMENT, ENABLED, SORT_ORDER"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "PRODUCT_APPKEY"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 存储第三方 API 密钥"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• PLATFORM='ALI_QWEN' 存储通义千问 API Key"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Oracle 元数据表"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• USER_TAB_COLUMNS：表列信息"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• USER_COL_COMMENTS：列注释"
    p.level = 1
    
    # ========== 第 7 页：API 接口 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "API 接口"
    tf = content.text_frame
    tf.text = "请求/响应格式"
    
    p = tf.add_paragraph()
    p.text = "请求"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "POST /api/nl-query/query"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Body: question=今天销售额是多少？"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "成功响应"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "{ success: true, sql: \"...\", data: [...], rowCount: N }"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "失败响应"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "{ success: false, message: \"...\" }"
    p.level = 1
    
    # ========== 第 8 页：技术栈 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "技术栈"
    tf = content.text_frame
    tf.text = "前端技术"
    
    p = tf.add_paragraph()
    p.text = "Vue 3.4.21 + Vite 5.2.6"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Element Plus 2.6.1"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Axios 1.6.8"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "后端技术"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Spring Boot 2.7.18"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Oracle JDBC 21.9.0.0"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "FastJSON2 2.0.43"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "OkHttp3（AI API 调用）"
    p.level = 1
    
    # ========== 结束页 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "谢谢"
    subtitle.text = "Q & A\n\n文档位置：/home/admin/.openclaw/workspace/智问核心流程泳道图_实际实现版.md"
    
    # 保存 PPT
    prs.save('/home/admin/.openclaw/workspace/智问核心流程.pptx')
    print("PPT 已生成：/home/admin/.openclaw/workspace/智问核心流程.pptx")

if __name__ == "__main__":
    create_presentation()
