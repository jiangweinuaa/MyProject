#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
智问核心流程 PPT v3 - 包含图形化泳道图
"""

from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()
    
    # ========== 封面页 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "智问系统核心流程泳道图"
    subtitle.text = "自然语言查询系统技术架构\n2026-04-14"
    
    # ========== 第 1 页：泳道图（图形） ==========
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # 空白布局
    title = slide.shapes.title
    title.text = "智问核心流程泳道图"
    
    # 添加泳道图图片
    left = Inches(0.5)
    top = Inches(1.5)
    slide.shapes.add_picture('/home/admin/.openclaw/workspace/智问泳道图.png', left, top, width=Inches(9))
    
    # ========== 第 2 页：泳道说明 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "6 个泳道职责"
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "1. 用户：输入问题、查看结果"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "2. 前端 (SmartQuery.vue)：界面展示、HTTP 请求、结果渲染"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. Controller (NLQueryController)：接收 HTTP 请求、参数校验"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "4. Service (NLQueryService)：SQL 生成、安全验证、执行查询"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "5. AI API (通义千问)：根据 Prompt 生成 SQL"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "6. 数据库 (Oracle)：表结构元数据、业务数据存储"
    p.level = 1
    
    # ========== 第 3 页：14 步流程 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "14 步完整流程"
    tf = content.text_frame
    
    steps = [
        "① 用户输入问题",
        "② 前端显示「AI 思考中...」",
        "③ 前端发送 POST 请求",
        "④ Controller 调用 Service",
        "⑤ Service 调用 generateSQL()",
        "⑥ 查询 AI_TABLE_FILTER 表",
        "⑦ 查询 USER_TAB_COLUMNS 元数据",
        "⑧ 调用通义千问 API",
        "⑨ AI 返回生成的 SQL",
        "⑩ Service 执行 validateSQL() 验证",
        "⑪ Service 执行 SQL 查询",
        "⑫ Service 返回 JSON 响应",
        "⑬ 前端渲染表格和 SQL",
        "⑭ 用户查看结果"
    ]
    
    for i, step in enumerate(steps):
        p = tf.add_paragraph()
        p.text = step
        p.level = 1
    
    # ========== 第 4 页：安全机制 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "安全机制（三层验证）"
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "第一层：SQL 类型检查"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "✓ 必须以 SELECT 开头"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✗ 禁止 DROP/DELETE/UPDATE/INSERT/TRUNCATE"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "第二层：表白名单检查"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "✓ 从 AI_TABLE_FILTER 读取允许的表"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✓ SQL 中必须包含允许的表名"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "第三层：只读账号执行"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "✓ JdbcTemplate 使用只读数据库账号"
    p.level = 1
    
    # ========== 第 5 页：技术栈 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "技术栈"
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "前端"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Vue 3.4.21 + Vite 5.2.6 + Element Plus 2.6.1"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "后端"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Spring Boot 2.7.18 + Oracle JDBC + OkHttp3"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "AI"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "通义千问 qwen-plus (dashscope.aliyuncs.com)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "数据库"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Oracle 11g"
    p.level = 1
    
    # ========== 结束页 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "谢谢"
    subtitle.text = "Q & A\n\n文档：/home/admin/.openclaw/workspace/智问泳道图_修正版.md"
    
    # 保存 PPT
    prs.save('/home/admin/.openclaw/workspace/智问泳道图_图形版.pptx')
    print("PPT 已生成：/home/admin/.openclaw/workspace/智问泳道图_图形版.pptx")

if __name__ == "__main__":
    create_presentation()
