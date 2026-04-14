#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
智问核心流程 PPT 生成脚本 v2 - 包含泳道图
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # ========== 封面页 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "智问系统核心流程泳道图"
    subtitle.text = "自然语言查询系统技术架构\n2026-04-14"
    
    # ========== 第 1 页：什么是泳道图 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "什么是泳道图？"
    tf = content.text_frame
    tf.text = "泳道图 vs 普通流程图"
    
    p = tf.add_paragraph()
    p.text = "泳道图特点"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 按角色/系统分区（泳道）"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 每个步骤放在对应的泳道中"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 清晰展示「谁在什么时候做什么」"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 跨泳道连接展示交互"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "智问系统的泳道"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "用户 | 前端 | Controller | Service | AI API | 数据库"
    p.level = 1
    
    # ========== 第 2 页：泳道职责表 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "泳道职责"
    tf = content.text_frame
    tf.text = "6 个泳道的职责分工"
    
    p = tf.add_paragraph()
    p.text = "1. 用户：输入问题、查看结果"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "2. 前端 (SmartQuery.vue)：界面展示、HTTP 请求"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. Controller：接收请求、参数校验"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "4. Service：SQL 生成、安全验证、执行查询"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "5. AI API (通义千问)：生成 SQL"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "6. 数据库 (Oracle)：表结构、业务数据"
    p.level = 1
    
    # ========== 第 3 页：核心流程泳道图 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "核心流程泳道图"
    tf = content.text_frame
    tf.text = "跨泳道交互流程"
    
    p = tf.add_paragraph()
    p.text = "用户 → 前端：输入问题"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "前端 → Controller：POST /api/nl-query/query"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Controller → Service：query(question)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Service → 数据库：查询表结构"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Service → AI API：调用通义千问"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "AI API → Service：返回 SQL"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Service → Service：validateSQL 验证"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Service → 数据库：执行 SQL"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Service → 前端：返回 JSON"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "前端 → 用户：渲染结果"
    p.level = 1
    
    # ========== 第 4 页：时序图 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "时序图（11 步交互）"
    tf = content.text_frame
    tf.text = "完整的请求 - 响应周期"
    
    p = tf.add_paragraph()
    p.text = "1. 用户输入问题"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "2. 前端显示「AI 思考中...」"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. 前端发送 POST 请求"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "4. Controller 调用 Service"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "5. Service 查询数据库元数据"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "6. Service 调用 AI API 生成 SQL"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "7. Service 执行安全验证"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "8. Service 执行 SQL 查询"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "9. Service 返回 JSON 响应"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "10. 前端渲染表格和 SQL"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "11. 用户查看结果"
    p.level = 1
    
    # ========== 第 5 页：前端泳道 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "前端泳道 (SmartQuery.vue)"
    tf = content.text_frame
    tf.text = "前端职责"
    
    p = tf.add_paragraph()
    p.text = "界面展示"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 聊天头部（标题 + 副标题）"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 消息容器（用户/机器人消息）"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 输入框 + 发送按钮"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "HTTP 请求"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• POST /api/nl-query/query"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Content-Type: application/x-www-form-urlencoded"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Body: question=xxx"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "结果渲染"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Markdown 表格转 HTML"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• SQL 代码块高亮显示"
    p.level = 1
    
    # ========== 第 6 页：后端泳道 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "后端泳道"
    tf = content.text_frame
    tf.text = "Controller + Service"
    
    p = tf.add_paragraph()
    p.text = "NLQueryController"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• @PostMapping /api/nl-query/query"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 接收 question 参数"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 调用 NLQueryService.query()"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "NLQueryService"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 调用 AISQLService.generateSQL()"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 调用 AISQLService.validateSQL()"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• JdbcTemplate 执行 SQL"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 构建 JSON 响应"
    p.level = 1
    
    # ========== 第 7 页：AI 集成泳道 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "AI 集成泳道"
    tf = content.text_frame
    tf.text = "通义千问 API 调用"
    
    p = tf.add_paragraph()
    p.text = "API 信息"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 端点：dashscope.aliyuncs.com"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 模型：qwen-plus"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• API Key：从 PRODUCT_APPKEY 表读取"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Prompt 构建"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 从 AI_TABLE_FILTER 获取表列表"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 从 USER_TAB_COLUMNS 获取列信息"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 从 USER_COL_COMMENTS 获取注释"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 9 条 SQL 生成规范"
    p.level = 1
    
    # ========== 第 8 页：安全验证泳道 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "安全验证泳道"
    tf = content.text_frame
    tf.text = "三层安全检查"
    
    p = tf.add_paragraph()
    p.text = "第一层：SQL 类型检查"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "✓ 必须以 SELECT 开头"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✗ 禁止 DROP/DELETE/UPDATE/INSERT"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "第二层：表白名单检查"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "✓ 从 AI_TABLE_FILTER 读取允许的表"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✓ SQL 中必须包含允许的表名"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "第三层：只读账号执行"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "✓ JdbcTemplate 使用只读数据库账号"
    p.level = 1
    
    # ========== 第 9 页：数据库泳道 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "数据库泳道"
    tf = content.text_frame
    tf.text = "Oracle 数据库表"
    
    p = tf.add_paragraph()
    p.text = "配置表"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• AI_TABLE_FILTER：AI 可访问的表白名单"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• PRODUCT_APPKEY：API Key 存储"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "元数据表（Oracle 系统表）"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• USER_TAB_COLUMNS：表列信息"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• USER_COL_COMMENTS：列注释"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "业务数据表"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• DCP_SALE：销售数据"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• DCP_STOCK：库存数据"
    p.level = 1
    
    # ========== 第 10 页：关键交互点 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "关键交互点"
    tf = content.text_frame
    tf.text = "9 个跨泳道交互"
    
    p = tf.add_paragraph()
    p.text = "1. 用户 → 前端：输入问题"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "2. 前端 → Controller：HTTP POST"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. Controller → Service：方法调用"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "4. Service → 数据库：查询元数据"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "5. Service → AI API：生成 SQL"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "6. Service → Service：安全验证"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "7. Service → 数据库：执行查询"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "8. Service → 前端：JSON 响应"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "9. 前端 → 用户：结果展示"
    p.level = 1
    
    # ========== 结束页 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "谢谢"
    subtitle.text = "Q & A\n\n文档位置：/home/admin/.openclaw/workspace/智问泳道图_修正版.md"
    
    # 保存 PPT
    prs.save('/home/admin/.openclaw/workspace/智问核心流程泳道图.pptx')
    print("PPT 已生成：/home/admin/.openclaw/workspace/智问核心流程泳道图.pptx")

if __name__ == "__main__":
    create_presentation()
